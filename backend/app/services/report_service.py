import os, json, re

def _safe_filename(name):
    return re.sub(r'[^\w\-_. ]', '', name).strip()[:50] or 'report'
from datetime import datetime
from typing import List, Optional
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession
from fastapi.encoders import jsonable_encoder

from ..models.report import Report
from ..models.biomechanical_data import BiomechanicalTestData
from ..models.project import BiomechanicalTestProject
from ..schemas.report import ReportCreate, ReportResponse
from ..core.config import STORAGE_DIR


def _generate_pptx(report_name, project_info, analysis_results):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = report_name; p.font.size = Pt(36); p.font.bold = True; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}"; p2.font.size = Pt(16); p2.alignment = PP_ALIGN.CENTER

    # Project info
    if project_info:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(6))
        tf = txBox.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = "Test Background"; p.font.size = Pt(28); p.font.bold = True
        for k, v in project_info.items():
            if v:
                p2 = tf.add_paragraph(); p2.text = f"{k}: {v}"; p2.font.size = Pt(14)

    # Analysis slides
    for i, ar in enumerate(analysis_results or []):
        if not isinstance(ar, dict): continue
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(6))
        tf = txBox.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.text = ar.get("metric_name", f"Analysis #{i+1}"); p.font.size = Pt(24); p.font.bold = True
        stats = ar.get("stats") or ar.get("result_data", {}).get("stats", {})
        for k in ["n", "mean", "std", "min", "max", "median"]:
            if k in stats:
                p2 = tf.add_paragraph(); p2.text = f"{k}: {stats[k]}"; p2.font.size = Pt(14)
        comparison = ar.get("comparison") or ar.get("result_data", {}).get("comparison", {})
        if comparison:
            p2 = tf.add_paragraph(); p2.text = "Group Comparison:"; p2.font.size = Pt(16); p2.font.bold = True
            for g, gs in comparison.items():
                p3 = tf.add_paragraph(); p3.text = f"  {g}: Mean={gs.get('mean','')}, SD={gs.get('std','')}"; p3.font.size = Pt(14)

    out = os.path.join(STORAGE_DIR, "reports", f"{_safe_filename(report_name)}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx")
    os.makedirs(os.path.dirname(out), exist_ok=True); prs.save(out); return out


def _generate_pdf(report_name, project_info, analysis_results):
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table
    from reportlab.lib.styles import getSampleStyleSheet

    out = os.path.join(STORAGE_DIR, "reports", f"{_safe_filename(report_name)}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf")
    os.makedirs(os.path.dirname(out), exist_ok=True)
    doc = SimpleDocTemplate(out, pagesize=A4)
    styles = getSampleStyleSheet()
    story = [Paragraph(f"<b>{report_name}</b>", styles["Title"]), Spacer(1, 12)]
    if project_info:
        story.append(Paragraph("<b>Test Background</b>", styles["Heading2"]))
        for k, v in project_info.items():
            if v: story.append(Paragraph(f"<b>{k}:</b> {v}", styles["Normal"]))
        story.append(Spacer(1, 12))
    for i, ar in enumerate(analysis_results or []):
        if not isinstance(ar, dict): continue
        story.append(Paragraph(f"<b>{ar.get('metric_name', f'Analysis #{i+1}')}</b>", styles["Heading2"]))
        stats = ar.get("stats") or ar.get("result_data", {}).get("stats", {})
        if stats: story.append(Table([[k, str(v)] for k, v in stats.items()]))
        comparison = ar.get("comparison") or ar.get("result_data", {}).get("comparison", {})
        if comparison:
            story.append(Paragraph("<b>Groups:</b>", styles["Heading3"]))
            for g, gs in comparison.items():
                story.append(Paragraph(f"{g}: Mean={gs.get('mean','')}, SD={gs.get('std','')}", styles["Normal"]))
        story.append(Spacer(1, 12))
    doc.build(story); return out


async def create_report(db, data: ReportCreate):
    r = Report(**data.model_dump()); db.add(r); await db.flush(); await db.refresh(r); return r

async def generate_report_file(db, report_id: int):
    result = await db.execute(select(Report).where(Report.id == report_id))
    r = result.scalar_one_or_none()
    if not r: return None
    p_result = await db.execute(select(BiomechanicalTestProject).where(BiomechanicalTestProject.id == r.project_id))
    p = p_result.scalar_one_or_none()
    proj_info = {"Project": p.project_no, "Name": p.project_name, "Brand": p.brand_name, "Shoe": p.shoe_name} if p else {}
    from ..models.analysis import Analysis
    ars = []
    aids = json.loads(json.dumps(r.analysis_ids)) if r.analysis_ids else []
    if isinstance(aids, list):
        for aid in aids:
            a_r = await db.execute(select(Analysis).where(Analysis.id == aid))
            a = a_r.scalar_one_or_none()
            if a: ars.append({"metric_name": a.analysis_name, "result_data": a.result_data})
    fp = _generate_pptx(r.report_name, proj_info, ars) if r.report_type == "PPTX" else _generate_pdf(r.report_name, proj_info, ars)
    from ..models.test_file import TestFile
    f = TestFile(file_no=f"RPT-{r.id}", file_name=os.path.basename(fp), file_type="PDF" if r.report_type == "PDF" else "PPT", storage_path=os.path.relpath(fp, STORAGE_DIR).replace("\\", "/"), file_size_bytes=os.path.getsize(fp), project_id=r.project_id, uploader="system", upload_time=datetime.now())
    db.add(f); await db.flush(); await db.refresh(f)
    r.file_id = f.id; r.status = "Generated"; await db.flush(); await db.refresh(r); return r

async def get_reports_list(db, project_id=None):
    q = select(Report)
    if project_id: q = q.where(Report.project_id == project_id)
    q = q.order_by(Report.created_at.desc()).limit(50)
    result = await db.execute(q)
    return jsonable_encoder([ReportResponse.model_validate(r).model_dump() for r in result.scalars().all()])

async def get_report(db, rid: int):
    r = await db.execute(select(Report).where(Report.id == rid)); return r.scalar_one_or_none()

async def delete_report(db, rid: int) -> bool:
    r = await get_report(db, rid)
    if not r: return False
    await db.delete(r); await db.flush(); return True
